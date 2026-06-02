import os
import re
from PyPDF2 import PdfReader, PdfWriter
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib import colors
from reportlab.lib.colors import Color, HexColor
from reportlab.lib.pagesizes import landscape, A4
from PIL import Image
import io
import fitz  # PyMuPDF - render each reportlab page straight to an image (no PyPDF2 merge)
from ai_vision import explain_image

font_heading = "Helvetica-Bold"
font_description = "Helvetica"
background_color = "#E7EDF1"

def clean_title(filename):
    """Clean the image title by removing numbers, underscores, and file extension."""
    # Remove file extension
    title = os.path.splitext(filename)[0]
    # Remove numbers and underscores from the beginning
    title = re.sub(r'^\d+[_-]*', '', title)
    # Replace remaining underscores with spaces and capitalize
    title = title.replace('_', ' ').strip().title()
    return title

def get_image_description(image_path):
    """Generate a one-line image description using Kimi AI."""
    try:
        return explain_image(image_path, "Generate a one-line explanation of this image.")
    except Exception as e:
        print(f"Error getting image description: {e}")
        return ""

def process_pdf(template_pdf_path, images_folder, output_pdf_path, start_page):
    """Process PDF and add images starting from specified page."""
    # Read existing PDF
    reader = PdfReader(template_pdf_path)
    writer = PdfWriter()

    # Copy pages up to start_page
    for i in range(start_page):
        writer.add_page(reader.pages[i])

    # Get sorted list of images
    images = sorted([f for f in os.listdir(images_folder)
                    if f.lower().endswith(('.png', '.jpg', '.jpeg'))])

    # Get template page for dimensions
    template_page = reader.pages[0]

    # Process each image
    for idx, image in enumerate(images, start=1):  # Start counting from 1
        image_path = os.path.join(images_folder, image)

        # Get image description
        description = get_image_description(image_path)

        # Create a new PDF in memory
        packet = io.BytesIO()
        # Pass the image index to determine width
        add_image_to_pdf(packet, image_path, description, template_page, idx)

        # Move to the beginning of the BytesIO buffer
        packet.seek(0)
        new_pdf = PdfReader(packet)

        # Add the new page
        writer.add_page(new_pdf.pages[0])

    # Copy remaining pages from template if any
    for i in range(start_page, len(reader.pages)):
        writer.add_page(reader.pages[i])

    # Write the final PDF
    with open(output_pdf_path, 'wb') as output_file:
        writer.write(output_file)

def add_image_to_pdf(packet, image_path, description, template_page, image_index):
    """Add image and description to a PDF page matching template dimensions."""
    # Get dimensions from template page
    template_width = float(template_page.mediabox.width)
    template_height = float(template_page.mediabox.height)

    # Create canvas with template dimensions
    can = canvas.Canvas(packet, pagesize=(template_width, template_height))

    # Add dotted background
    dot_radius = 1  # Radius of each dot
    dot_spacing = 20  # Space between dots (both horizontally and vertically)
    dot_color = Color(0.8, 0.8, 0.8, alpha=1)  # Light gray color for dots

    for x in range(0, int(template_width), dot_spacing):
        for y in range(0, int(template_height), dot_spacing):
            can.setFillColor(dot_color)
            can.circle(x, y, dot_radius, fill=True, stroke=False)

    # Add gradient background
    can.setFillColor(HexColor('#E7EDF1'))  # Light gray background
    can.rect(0, 0, template_width, template_height, fill=True)

    # DigiRocket brand stamp (top-right of every data slide). The template
    # cover/agenda/thank-you pages already carry the logo, so this only runs
    # for the data slides built by this function. Acts as a visual mark of
    # authenticity on every page.
    _brand_text = "DigiRocket"
    _brand_font = "Helvetica-Bold"
    _brand_size = 18
    can.setFont(_brand_font, _brand_size)
    can.setFillColor(HexColor('#002060'))
    _brand_w = can.stringWidth(_brand_text, _brand_font, _brand_size)
    can.drawString(template_width - _brand_w - 30, template_height - 45, _brand_text)
    # Brand-coloured accent pill under the wordmark
    can.setFillColor(HexColor('#C9F31D'))
    can.rect(template_width - _brand_w - 30, template_height - 55, _brand_w, 4, fill=True, stroke=False)

    # Load and resize image while maintaining aspect ratio
    img = Image.open(image_path)
    img_width, img_height = img.size
    aspect = img_width / float(img_height)

    # Calculate new dimensions
    if image_index == 1 or image_index == 7:
        # Keep original size for 1st and 7th images
        new_height = template_height * 0.5
        new_width = new_height * aspect
        max_width_ratio = 0.9  # Constrain width to 90% of page width
    else:
        # Make other images larger and wider
        new_height = template_height * 0.6  # Increase height to 70% of page height
        new_width = new_height * aspect
        max_width_ratio = 1.0  # Allow width to be up to 90% of page width

    # Constrain width if it exceeds the maximum allowed width
    if new_width > template_width * max_width_ratio:
        new_width = template_width * max_width_ratio
        new_height = new_width / aspect  # Adjust height to maintain aspect ratio

    # Center the image
    x = (template_width - new_width) / 2
    y = (template_height - new_height) / 2

    # Draw image
    can.drawImage(image_path, x, y, width=new_width, height=new_height)

    # Add rounded rectangle border around the image
    border_thickness = 10  # Thickness of the border
    border_color = HexColor('#002060')  # Dark blue color for border
    corner_radius = 20  # Rounded corners

    # Draw the rounded rectangle border
    can.setStrokeColor(border_color)
    can.setLineWidth(border_thickness)
    can.roundRect(
        x - border_thickness / 2,
        y - border_thickness / 2,
        new_width + border_thickness,
        new_height + border_thickness,
        corner_radius,
        stroke=True,
        fill=False
    )

    # Get clean image title
    image_title = clean_title(os.path.basename(image_path))

    # Add title text above the image with gradient effect
    can.setFont("Helvetica-Bold", 40)
    title_width = can.stringWidth(image_title, "Helvetica-Bold", 40)
    title_x = (template_width - title_width) / 2
    title_y = y + new_height + 100

    # Gradient fill for title
    can.setFillColor(HexColor('#002060'))  # Start color
    can.drawString(title_x, title_y, image_title)

    # Add shadow effect for title
    can.setFillColor(colors.black)
    can.drawString(title_x + 2, title_y - 2, image_title)  # Shadow offset

    # Reset fill color back to black for other text
    can.setFillColor(colors.black)

    # Add description text below the image
    can.setFont("Helvetica", 18)
    description_lines = []

    # Split description into multiple lines if too long
    max_width = template_width * 0.8
    while can.stringWidth(description, "Helvetica", 18) > max_width:
        words = description.split()
        current_line = []
        while words and can.stringWidth(' '.join(current_line + [words[0]]), "Helvetica", 18) <= max_width:
            current_line.append(words.pop(0))
        description_lines.append(' '.join(current_line))
        description = ' '.join(words)
    if description:
        description_lines.append(description)

    # Draw description lines with padding and spacing
    text_y = y - 100
    for line in description_lines:
        text_width = can.stringWidth(line, "Helvetica", 18)
        text_x = (template_width - text_width) / 2
        can.drawString(text_x, text_y, line)
        text_y -= 30  # Increased spacing between lines

    can.save()


def build_slide_images(template_pdf_path, images_folder, output_dir, start_page,
                       zoom=None, ai_text='', tables_by_image=None):
    """Build the slide images for the video WITHOUT merging PDFs.

    The old approach merged many reportlab pages into one PDF with PyPDF2, which
    dropped image XObjects on some pages (-> blank slides + 'cannot find XObject'
    MuPDF errors). Here each reportlab page is rendered straight to an image with
    PyMuPDF, so nothing is merged and no image is ever lost. Also faster.

    Output: page_1.png, page_2.png, ... in `output_dir` (HD via `zoom`).

    Also writes `manifest.json` describing every slide so a downstream PPT
    builder can place the ORIGINAL screenshot + the AI description as separate
    editable text boxes (and template covers as static branded images).

    `ai_text` (optional): the full AI Strategy Memo text. If supplied AND the
    captured screenshots include the AI Insights slide (filename containing
    "ai-insights"), the manifest marks that slide as type='ai-text' and stores
    the full memo so the PPT builder can render it as editable text boxes
    instead of an image.
    """
    import json
    os.makedirs(output_dir, exist_ok=True)
    # Slide-render DPI. Default 1.5 so a 22-slide PDF fits comfortably under
    # Render's 512 MB free-tier RAM. Each step up of `zoom` quadruples
    # the pixmap size; 1.5 -> ~1440x810 per slide which is sharper than the
    # email/PDF view ever shows. Bump via PDF_RENDER_ZOOM if you're on a
    # bigger box and want truly print-grade output.
    if zoom is None:
        try:
            zoom = float(os.getenv('PDF_RENDER_ZOOM', '1.5'))
        except ValueError:
            zoom = 1.5
    template = fitz.open(template_pdf_path)
    template_page_ref = PdfReader(template_pdf_path).pages[0]  # only for dimensions
    matrix = fitz.Matrix(zoom, zoom)
    page_num = 0
    manifest = {'slides': []}

    def save_pixmap(pix):
        nonlocal page_num
        page_num += 1
        pix.save(os.path.join(output_dir, f"page_{page_num}.png"))

    # 1) Template cover pages (rendered directly from the template - no merge).
    # SLIDE 1 + 2 OVERRIDE: if a custom image exists in static/images/, use it
    # instead of the corresponding Template.pdf page.
    #   - cover.png   -> slide 1 (front cover)
    #   - agenda.png  -> slide 2 (agenda)
    # The closing Thank-You slide still comes from Template.pdf.
    # Accepts .png / .jpg / .jpeg and the common Windows ".png.jpg" variant
    # (Windows often appends the real extension when you save a JPG with a
    # ".png" filename in Save-As).
    script_dir = os.path.dirname(os.path.abspath(__file__))
    images_root = os.path.join(script_dir, 'static', 'images')

    def _find_override(stem):
        # Windows Save-As frequently appends the real extension, so a file the
        # user "named" stem.png often lands on disk as stem.png.png or
        # stem.png.jpg. Accept every common variant.
        for ext in ('.png', '.jpg', '.jpeg',
                    '.png.png', '.png.jpg', '.png.jpeg',
                    '.jpg.png', '.jpg.jpg'):
            p = os.path.join(images_root, stem + ext)
            if os.path.exists(p):
                return p
        return None

    custom_cover = _find_override('cover')
    custom_agenda = _find_override('agenda')
    overrides = {0: custom_cover, 1: custom_agenda}

    def _render_override_to_pixmap(image_path):
        """Image -> single-page PDF -> rasterised pixmap at slide zoom."""
        cover_doc = fitz.open(image_path)
        pdfbytes = cover_doc.convert_to_pdf()
        cover_doc.close()
        tmp = fitz.open(stream=pdfbytes, filetype='pdf')
        try:
            return tmp[0].get_pixmap(matrix=matrix)
        finally:
            tmp.close()

    cover_count = min(start_page, len(template))
    for i in range(cover_count):
        override_path = overrides.get(i)
        if override_path:
            try:
                save_pixmap(_render_override_to_pixmap(override_path))
            except Exception as e:
                print(f"  (custom slide {i+1} failed, falling back to template: {e})")
                save_pixmap(template[i].get_pixmap(matrix=matrix))
        else:
            save_pixmap(template[i].get_pixmap(matrix=matrix))
        manifest['slides'].append({
            'type': 'template',
            'slide_image': f'page_{page_num}.png',
        })

    # 2) Data slides: build each reportlab page, render it straight to an image
    images = sorted([f for f in os.listdir(images_folder)
                     if f.lower().endswith(('.png', '.jpg', '.jpeg'))])

    # Per-image AI captions (one-line Kimi vision description under each
    # slide) are OFF by default. They were the single biggest cause of the
    # Render free-tier 502 / OOM: every call loads + downscales + base64s
    # the full screenshot in Python memory, and at 20 slides that easily
    # tipped the 512 MB worker over.
    #
    # Set PDF_AI_CAPTIONS=1 in the environment to turn them back on (e.g.
    # on a paid Render tier with more RAM). The Summary & Conclusion AI
    # memo on the report is unaffected - it runs client-side from the
    # view, not here.
    from concurrent.futures import ThreadPoolExecutor as _TPE

    _ai_captions_on = os.getenv('PDF_AI_CAPTIONS', '0').lower() in ('1', 'true', 'yes', 'on')

    def _desc_for(image_filename):
        try:
            return get_image_description(os.path.join(images_folder, image_filename))
        except Exception as e:
            print(f"  (description failed for {image_filename}: {e})")
            return ""

    if images and _ai_captions_on:
        # Pool cap of 4: each worker holds a full image in memory.
        with _TPE(max_workers=min(4, len(images))) as _ex:
            _descriptions = list(_ex.map(_desc_for, images))
    else:
        # Empty captions -> add_image_to_pdf just leaves the description
        # area blank, slides still get the image + title + brand stamp.
        _descriptions = ['' for _ in images]

    # Pre-render each data slide's ReportLab PDF page bytes in parallel.
    # add_image_to_pdf is CPU-bound (PIL resize, reportlab canvas) and
    # reads filesystem images, so threads give a real speedup. We still
    # rasterise+save serially below (PyMuPDF docs are single-thread per doc).
    def _build_packet(args):
        idx_, image_filename, desc_ = args
        ipath = os.path.join(images_folder, image_filename)
        pkt = io.BytesIO()
        try:
            add_image_to_pdf(pkt, ipath, desc_, template_page_ref, idx_)
        except Exception as e:
            print(f"  (slide build failed for {image_filename}: {e})")
            return None
        return pkt.getvalue()

    _packet_args = [(i + 1, name, _descriptions[i] if i < len(_descriptions) else "")
                    for i, name in enumerate(images)]
    _packets = []
    if _packet_args:
        # Pool cap matches the AI caption pool - each worker opens the
        # source screenshot in Pillow + builds a ReportLab page, which
        # peaks at a few MB. Override with PDF_RENDER_WORKERS env var
        # on bigger boxes.
        _render_workers = int(os.getenv('PDF_RENDER_WORKERS', '4'))
        with _TPE(max_workers=min(_render_workers, len(_packet_args))) as _ex:
            _packets = list(_ex.map(_build_packet, _packet_args))

    for idx, image in enumerate(images, start=1):
        image_path = os.path.join(images_folder, image)
        description = _descriptions[idx - 1] if idx - 1 < len(_descriptions) else ""
        title = clean_title(os.path.basename(image_path))
        pkt_bytes = _packets[idx - 1] if idx - 1 < len(_packets) else None
        if not pkt_bytes:
            continue
        single = fitz.open(stream=pkt_bytes, filetype='pdf')
        save_pixmap(single[0].get_pixmap(matrix=matrix))
        single.close()
        slide_entry = {
            'type': 'data',
            'slide_image': f'page_{page_num}.png',
            # Store the ORIGINAL screenshot path so the editable PPT can drop
            # it in as a movable, non-edited image (preserves data authenticity).
            'original_image': image_path,
            'title': title,
            'description': description,
        }
        # If the frontend extracted a real <table> for this screenshot, attach
        # headers+rows to the manifest entry. The PPT builder will then render
        # it as a NATIVE editable PowerPoint table instead of just an image.
        if tables_by_image:
            tbl = tables_by_image.get(os.path.basename(image_path))
            if tbl and (tbl.get('headers') or tbl.get('rows')):
                slide_entry['table'] = tbl
        # If this is the closing Summary & Conclusion slide AND we received
        # the full memo text, tag it so the PPT / PDF builders render it as
        # editable text instead of a flattened image.
        # (Matches both the new filename `summary-and-conclusion` and the
        # legacy `ai-insights` names so older AIVideo folders still work.)
        lower_name = os.path.basename(image_path).lower()
        if ai_text and (
            'summary-and-conclusion' in lower_name
            or 'summary_and_conclusion' in lower_name
            or 'ai-insights' in lower_name
            or 'ai_insights' in lower_name
        ):
            slide_entry['type'] = 'ai-text'
            slide_entry['ai_text'] = ai_text
            # Override the cleaned-from-filename title so the slide reads as
            # a neutral "Summary & Conclusion" in PDF and PPT (the client
            # shouldn't see "AI" in the heading; the disclaimer below the
            # heading already covers that the content was AI-assisted).
            slide_entry['title'] = 'Summary & Conclusion'
        manifest['slides'].append(slide_entry)

    # 3) Remaining template pages (Thank You etc.) — with optional
    # thankyou.png override for the FIRST remaining page (the closing
    # Thank You slide). Same pattern as cover / agenda overrides above.
    custom_thankyou = _find_override('thankyou')
    for offset, i in enumerate(range(start_page, len(template))):
        if offset == 0 and custom_thankyou:
            try:
                save_pixmap(_render_override_to_pixmap(custom_thankyou))
            except Exception as e:
                print(f"  (custom thankyou failed, falling back to template: {e})")
                save_pixmap(template[i].get_pixmap(matrix=matrix))
        else:
            save_pixmap(template[i].get_pixmap(matrix=matrix))
        manifest['slides'].append({
            'type': 'template',
            'slide_image': f'page_{page_num}.png',
        })

    template.close()

    # Manifest lets /download-report-pptx rebuild the deck with editable title +
    # description text boxes around the original screenshots, without having to
    # re-call the AI or re-scrape the dashboard.
    try:
        with open(os.path.join(output_dir, 'manifest.json'), 'w', encoding='utf-8') as f:
            json.dump(manifest, f, indent=2, ensure_ascii=False)
    except Exception as e:
        print(f"  (could not save manifest: {e})")

    print(f"Built {page_num} slide images in {output_dir} (no-merge, HD)")
    return page_num


def build_pdf_from_images(images_dir, output_pdf):
    """Combine the slide images (page_1.png ...) into one downloadable PDF report.

    Uses PyMuPDF to insert each image as a page - reliable, no XObject merge bug.

    Special case: if `manifest.json` marks a slide as type='ai-text', that slide
    is rendered as a REAL TEXT page (selectable / copy-able) instead of a
    flattened image, so the AI Strategy Memo stays searchable and editable.
    """
    import json

    def _page_num(name):
        m = re.search(r'page_(\d+)', name)
        return int(m.group(1)) if m else 0

    # Load the manifest so we know which slide is the AI text slide.
    manifest_path = os.path.join(images_dir, 'manifest.json')
    manifest = {}
    if os.path.exists(manifest_path):
        try:
            with open(manifest_path, 'r', encoding='utf-8') as f:
                manifest = json.load(f)
        except Exception as e:
            print(f"  (could not load manifest: {e})")

    # Map slide_image filename -> slide entry, so we can look up type/ai_text
    slide_by_image = {}
    for s in (manifest.get('slides') or []):
        si = s.get('slide_image')
        if si:
            slide_by_image[si] = s

    images = sorted(
        [f for f in os.listdir(images_dir) if f.lower().endswith(('.png', '.jpg', '.jpeg'))],
        key=_page_num,
    )

    # Standard slide dimensions (matches the screenshot canvas: 16:9 widescreen
    # at 72 dpi PostScript points -> 960 x 540 pt is a clean 16:9 ratio).
    PAGE_W, PAGE_H = 960.0, 540.0

    def _is_heading(line):
        t = line.strip()
        if not t.endswith(':') or len(t) < 4:
            return False
        body = t[:-1]
        letters = sum(1 for c in body if c.isalpha())
        if not letters:
            return False
        upper = sum(1 for c in body if c.isupper())
        return (upper / letters) >= 0.6

    def _render_ai_text_page(out_doc, title_text, memo_text):
        """Insert a PyMuPDF text-rendered page (selectable text, no image)."""
        page = out_doc.new_page(width=PAGE_W, height=PAGE_H)
        # Light grey background to match the rest of the report
        page.draw_rect(fitz.Rect(0, 0, PAGE_W, PAGE_H), color=None,
                       fill=(0.906, 0.929, 0.945), overlay=False)
        # DigiRocket brand stamp top-right
        page.insert_text(fitz.Point(PAGE_W - 110, 26), 'DigiRocket',
                         fontname='hebo', fontsize=11, color=(0.0, 0.13, 0.38))
        page.draw_line(fitz.Point(PAGE_W - 110, 30), fitz.Point(PAGE_W - 30, 30),
                       color=(0.788, 0.953, 0.114), width=2)
        # Title (centred)
        title_y = 70
        title_size = 20
        title_w = fitz.get_text_length(title_text, fontname='hebo', fontsize=title_size)
        page.insert_text(fitz.Point((PAGE_W - title_w) / 2, title_y),
                         title_text, fontname='hebo', fontsize=title_size,
                         color=(0.0, 0.13, 0.38))
        # Body: paragraph-per-line, wrapped to page width
        body_left = 50
        body_right = PAGE_W - 50
        body_width = body_right - body_left
        cursor_y = title_y + 30
        lines = (memo_text or '').split('\n')
        for raw in lines:
            line = raw.rstrip()
            if not line.strip():
                cursor_y += 6
                continue
            is_head = _is_heading(line)
            font = 'hebo' if is_head else 'helv'
            size = 12 if is_head else 10
            color = (0.0, 0.13, 0.38) if is_head else (0.122, 0.161, 0.216)
            text = line.strip().rstrip(':') if is_head else line
            # word-wrap manually using get_text_length
            words = text.split(' ')
            cur = ''
            while words:
                w = words.pop(0)
                trial = (cur + ' ' + w).strip() if cur else w
                if fitz.get_text_length(trial, fontname=font, fontsize=size) <= body_width:
                    cur = trial
                else:
                    if cur:
                        page.insert_text(fitz.Point(body_left, cursor_y), cur,
                                         fontname=font, fontsize=size, color=color)
                        cursor_y += size + 4
                    cur = w
                if cursor_y > PAGE_H - 30:
                    break
            if cur and cursor_y <= PAGE_H - 30:
                page.insert_text(fitz.Point(body_left, cursor_y), cur,
                                 fontname=font, fontsize=size, color=color)
                cursor_y += size + 6
            if is_head:
                cursor_y += 2
            if cursor_y > PAGE_H - 30:
                break

    # Parallel pre-render: each screenshot -> single-page PDF bytes runs in
    # parallel threads (PyMuPDF + PIL release the GIL during image work, so
    # threads give a real ~3-4x speedup). The final assembly then inserts
    # the pre-rendered pages in order — assembly is fast but MUST be serial
    # (fitz docs are not thread-safe for cross-doc inserts).
    from concurrent.futures import ThreadPoolExecutor

    def _img_to_pdf_bytes(name):
        try:
            imgdoc = fitz.open(os.path.join(images_dir, name))
            try:
                return imgdoc.convert_to_pdf()
            finally:
                imgdoc.close()
        except Exception as e:
            print(f"  (image->pdf failed for {name}: {e})")
            return None

    # Only pre-render the slides that will be inserted as images (skip the
    # AI text slide, which is built directly in the assembly loop).
    image_names = [n for n in images
                   if (slide_by_image.get(n) or {}).get('type') != 'ai-text'
                   or not (slide_by_image.get(n) or {}).get('ai_text')]
    rendered = {}
    if image_names:
        with ThreadPoolExecutor(max_workers=min(8, len(image_names))) as ex:
            futures = {ex.submit(_img_to_pdf_bytes, n): n for n in image_names}
            for f in futures:
                rendered[futures[f]] = f.result()

    doc = fitz.open()
    for name in images:
        slide_meta = slide_by_image.get(name) or {}
        if slide_meta.get('type') == 'ai-text' and slide_meta.get('ai_text'):
            _render_ai_text_page(
                doc,
                title_text=slide_meta.get('title') or 'Summary & Conclusion',
                memo_text=slide_meta.get('ai_text', ''),
            )
            continue
        pdfbytes = rendered.get(name)
        if not pdfbytes:
            continue
        imgpdf = fitz.open("pdf", pdfbytes)
        doc.insert_pdf(imgpdf)
        imgpdf.close()
    doc.save(output_pdf)
    doc.close()
    print(f"PDF report saved: {output_pdf} ({len(images)} pages)")
    return output_pdf


def build_pptx_from_images(images_dir, output_pptx):
    """Combine the slide images (page_1.png ...) into a downloadable PPTX deck.

    Each slide is one full-slide image (the same images used for the PDF), so the
    deck matches the report exactly. The PMO team can then open it in PowerPoint
    to reorder/remove slides, add their own slides, or add speaker notes.
    """
    from pptx import Presentation
    from pptx.util import Emu, Inches, Pt
    from pptx.dml.color import RGBColor

    def _page_num(name):
        m = re.search(r'page_(\d+)', name)
        return int(m.group(1)) if m else 0

    images = sorted(
        [f for f in os.listdir(images_dir) if f.lower().endswith(('.png', '.jpg', '.jpeg'))],
        key=_page_num,
    )

    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 13.333in -> 16:9 widescreen
    prs.slide_height = Emu(6858000)   # 7.5in
    blank_layout = prs.slide_layouts[6]
    sw, sh = int(prs.slide_width), int(prs.slide_height)
    slide_aspect = sw / float(sh)

    for name in images:
        path = os.path.join(images_dir, name)
        slide = prs.slides.add_slide(blank_layout)
        with Image.open(path) as im:
            iw, ih = im.size
        img_aspect = iw / float(ih)
        # Fit the image inside the slide, preserving aspect ratio, centered.
        if img_aspect > slide_aspect:
            w = sw
            h = int(sw / img_aspect)
        else:
            h = sh
            w = int(sh * img_aspect)
        left = int((sw - w) / 2)
        top = int((sh - h) / 2)
        slide.shapes.add_picture(path, left, top, width=w, height=h)

    prs.save(output_pptx)
    print(f"PPTX (image slides) saved: {output_pptx} ({len(images)} slides)")
    return output_pptx


def build_editable_image_pptx(aivideo_dir, output_pptx):
    """Build a hybrid PPTX:
      - Template slides (cover / agenda / thank-you) are full-slide images so
        the branded look is preserved exactly.
      - Data slides have THREE separate movable/editable elements:
            1. Title text box (top, editable)
            2. The original captured screenshot (centred, MOVABLE; the image
               itself is not editable so the data stays authentic).
            3. Description text box (bottom, editable — same AI-generated copy
               that appears under the image in the PDF).

    This keeps the visuals trustworthy (no faking numbers in a table) while
    still letting the PMO team rewrite the narrative around the screenshots.

    Reads ``manifest.json`` written by :func:`build_slide_images` so we know
    which slide is which and where the original screenshot lives.
    """
    import json
    from pptx import Presentation
    from pptx.util import Emu, Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    manifest_path = os.path.join(aivideo_dir, 'manifest.json')
    if not os.path.exists(manifest_path):
        # Older AIVideo folders (built before this code shipped) have no
        # manifest -> fall back to the simple image-only deck so the user
        # still gets a usable PPT.
        return build_pptx_from_images(aivideo_dir, output_pptx)

    with open(manifest_path, 'r', encoding='utf-8') as f:
        manifest = json.load(f)

    NAVY = RGBColor(0x00, 0x20, 0x60)
    DARK = RGBColor(0x1F, 0x29, 0x37)
    BRAND_LIME = RGBColor(0xC9, 0xF3, 0x1D)

    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 13.333in -> 16:9 widescreen
    prs.slide_height = Emu(6858000)   # 7.5in
    blank_layout = prs.slide_layouts[6]
    SW, SH = int(prs.slide_width), int(prs.slide_height)
    slide_aspect = SW / float(SH)

    def add_full_slide_image(slide, img_path):
        with Image.open(img_path) as im:
            iw, ih = im.size
        img_aspect = iw / float(ih)
        if img_aspect > slide_aspect:
            w = SW; h = int(SW / img_aspect)
        else:
            h = SH; w = int(SH * img_aspect)
        left = int((SW - w) / 2)
        top = int((SH - h) / 2)
        slide.shapes.add_picture(img_path, left, top, width=w, height=h)

    def add_brand_stamp(slide):
        """Drop a movable 'DigiRocket' wordmark + lime accent on the slide
        (top-right). Both shapes are independent, so the user can drag /
        resize / restyle them in PowerPoint, but they're there by default
        so every page carries the brand mark."""
        from pptx.enum.shapes import MSO_SHAPE
        # Tight text box: no extra padding so the lime accent sits flush.
        tb = slide.shapes.add_textbox(SW - Inches(2.2), Inches(0.22),
                                      Inches(2.0), Inches(0.32))
        tf = tb.text_frame
        tf.word_wrap = False
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        p = tf.paragraphs[0]
        p.text = "DigiRocket"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.alignment = PP_ALIGN.RIGHT
        # Thin lime accent immediately under the wordmark (2pt below the
        # text baseline). Right-edge aligns with the text box's right edge.
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        SW - Inches(1.2), Inches(0.52),
                                        Inches(1.0), Pt(2.5))
        accent.fill.solid()
        accent.fill.fore_color.rgb = BRAND_LIME
        accent.line.fill.background()

    for s in manifest.get('slides', []):
        slide = prs.slides.add_slide(blank_layout)
        slide_type = s.get('type', 'data')
        slide_image_path = os.path.join(aivideo_dir, s.get('slide_image', ''))

        if slide_type == 'template':
            # Branded cover / agenda / thank-you page -> keep as a static image.
            # Template pages already carry the DigiRocket logo, so we skip the
            # extra brand stamp to avoid duplicates.
            if os.path.exists(slide_image_path):
                add_full_slide_image(slide, slide_image_path)
            continue

        # Every non-template slide carries a movable "DigiRocket" wordmark top-right
        add_brand_stamp(slide)

        # ---- AI Strategy Memo slide: render as EDITABLE TEXT, not image -------
        # The PMO team needs to be able to tweak the AI's wording before
        # sharing the deck, so this slide is built from real text boxes
        # (title + body) instead of a flattened screenshot.
        if slide_type == 'ai-text':
            title_text = s.get('title', 'Summary & Conclusion')
            memo_text = (s.get('ai_text') or '').strip() or '(AI memo not available)'

            # Title at top
            tb_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3),
                                                SW - Inches(1), Inches(0.85))
            tf = tb_title.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = NAVY
            p.alignment = PP_ALIGN.CENTER

            # Body: one editable text frame containing the full memo. Each
            # logical line in the memo becomes a paragraph in the text frame.
            # UPPERCASE LINES ending with ':' are rendered as bold section
            # headers; bullet/numbered lines keep their leading marker so the
            # user sees the same structure as on the live page.
            tb_body = slide.shapes.add_textbox(Inches(0.5), Inches(1.3),
                                               SW - Inches(1), SH - Inches(1.8))
            bf = tb_body.text_frame
            bf.word_wrap = True
            try:
                bf.auto_size = None  # let the text overflow rather than shrink
            except Exception:
                pass

            def _is_heading(line):
                t = line.strip()
                if not t.endswith(':') or len(t) < 4:
                    return False
                body = t[:-1]
                letters = sum(1 for c in body if c.isalpha())
                if not letters:
                    return False
                upper = sum(1 for c in body if c.isupper())
                return (upper / letters) >= 0.6

            lines = memo_text.split('\n')
            first = True
            for raw in lines:
                line = raw.rstrip()
                if first:
                    para = bf.paragraphs[0]
                    first = False
                else:
                    para = bf.add_paragraph()
                if not line.strip():
                    # Blank spacer paragraph
                    para.text = ''
                    para.font.size = Pt(6)
                    continue
                if _is_heading(line):
                    para.text = line.strip().rstrip(':')
                    para.font.size = Pt(13)
                    para.font.bold = True
                    para.font.color.rgb = NAVY
                else:
                    para.text = line
                    para.font.size = Pt(11)
                    para.font.color.rgb = DARK
                    para.font.bold = False
            # Skip the image-based path for this slide
            continue
        # -----------------------------------------------------------------

        # ---- Data slide: title text + (native table OR image) + description -
        title = s.get('title', '') or ''
        desc = s.get('description', '') or ''
        orig_img = s.get('original_image', '') or ''
        table_data = s.get('table')  # {headers:[], rows:[[...]]} if available

        # Title text box (top) — editable in PowerPoint.
        tb_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3),
                                            SW - Inches(1), Inches(0.85))
        tf = tb_title.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.alignment = PP_ALIGN.CENTER

        rendered_as_table = False
        if table_data and (table_data.get('headers') or table_data.get('rows')):
            try:
                headers = table_data.get('headers') or []
                rows = table_data.get('rows') or []
                # Normalise: every row should have the same column count as headers
                ncols = max(len(headers), max((len(r) for r in rows), default=0))
                if not headers:
                    headers = [''] * ncols
                else:
                    headers = list(headers) + [''] * (ncols - len(headers))
                norm_rows = []
                for r in rows:
                    rr = list(r) + [''] * (ncols - len(r))
                    norm_rows.append(rr[:ncols])
                nrows = 1 + len(norm_rows)
                # Cap visible rows so the table fits — full data still in image
                # backup below if it overflows. 14 body rows + 1 header is a
                # comfortable max on a 16:9 slide.
                MAX_BODY_ROWS = 14
                if len(norm_rows) > MAX_BODY_ROWS:
                    norm_rows = norm_rows[:MAX_BODY_ROWS]
                    nrows = 1 + MAX_BODY_ROWS

                tbl_w = SW - Inches(1.0)
                tbl_left = Inches(0.5)
                tbl_top = Inches(1.3)
                # Height scales with row count; cap so it never collides with
                # the description box at the bottom.
                row_h_in = 0.32
                tbl_h_in = min(4.3, max(1.0, nrows * row_h_in))
                tbl_h = Inches(tbl_h_in)

                shape = slide.shapes.add_table(nrows, ncols, tbl_left, tbl_top, tbl_w, tbl_h)
                table = shape.table

                # Header row styling
                for ci, h in enumerate(headers):
                    cell = table.cell(0, ci)
                    cell.text = str(h)
                    for para in cell.text_frame.paragraphs:
                        para.alignment = PP_ALIGN.LEFT
                        for run in para.runs:
                            run.font.size = Pt(11)
                            run.font.bold = True
                            run.font.color.rgb = WHITE if False else RGBColor(0xFF, 0xFF, 0xFF)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = NAVY

                # Body rows
                for ri, row in enumerate(norm_rows, start=1):
                    for ci, val in enumerate(row):
                        cell = table.cell(ri, ci)
                        cell.text = str(val)
                        for para in cell.text_frame.paragraphs:
                            para.alignment = PP_ALIGN.LEFT
                            for run in para.runs:
                                run.font.size = Pt(10)
                                run.font.color.rgb = DARK
                        if ri % 2 == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
                        else:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

                rendered_as_table = True
            except Exception as e:
                print(f"  (native table render failed for {title}: {e}) — falling back to image")
                rendered_as_table = False

        if not rendered_as_table:
            # Original screenshot (centred, MOVABLE). The user can drag it but
            # the image content itself is locked = data stays trustworthy.
            img_to_place = orig_img if os.path.exists(orig_img) else slide_image_path
            if os.path.exists(img_to_place):
                with Image.open(img_to_place) as im:
                    iw, ih = im.size
                img_aspect = iw / float(ih)
                # Reserve: title 0.3-1.15, image area 1.3-5.6, description 5.7-7.2
                max_img_h = Inches(4.3)
                max_img_w = SW - Inches(1.4)
                h = max_img_h
                w = int(h * img_aspect)
                if w > max_img_w:
                    w = max_img_w
                    h = int(w / img_aspect)
                left = int((SW - w) / 2)
                top = int(Inches(1.3) + (max_img_h - h) / 2)
                slide.shapes.add_picture(img_to_place, left, top, width=w, height=h)

        # Description text box (bottom) — editable, pre-filled with the same
        # AI-generated explanation that appears under the image in the PDF.
        tb_desc = slide.shapes.add_textbox(Inches(0.5), Inches(5.75),
                                           SW - Inches(1), Inches(1.5))
        tf = tb_desc.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = DARK
        p.alignment = PP_ALIGN.CENTER

    prs.save(output_pptx)
    print(f"Editable PPTX (movable images + editable text) saved: {output_pptx}")
    return output_pptx


def build_editable_pptx(context, output_pptx, slide_images=None):
    """Build a PowerPoint report.

    If `slide_images` (list of page_N.png paths) is given, those are added first
    as full-slide images so the deck VISUALLY matches the PDF (cover, colours,
    charts, screenshots). Then native EDITABLE slides (overview metrics + data
    tables) are appended so the PMO team can still edit the numbers/text.

    context = {
        'title': str, 'subtitle': str,
        'metrics': [(label, value), ...],
        'tables': [{'title': str, 'headers': [...], 'rows': [[...], ...]}, ...],
    }
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    NAVY = RGBColor(0x00, 0x20, 0x60)
    LIME = RGBColor(0x6F, 0xA8, 0x1F)
    DARK = RGBColor(0x1F, 0x29, 0x37)
    GREY = RGBColor(0x6B, 0x72, 0x80)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT = RGBColor(0xF1, 0xF5, 0xF9)

    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 13.333in -> 16:9
    prs.slide_height = Emu(6858000)   # 7.5in
    blank = prs.slide_layouts[6]
    SW, SH = int(prs.slide_width), int(prs.slide_height)

    def title_bar(slide, text):
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), SW - Inches(1), Inches(0.9))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = NAVY

    # ---- If slide images are given, the deck is JUST those (matches the PDF
    #      exactly). No duplicated "Editable Data" tables. ----
    if slide_images:
        slide_aspect = SW / float(SH)
        for path in slide_images:
            if not os.path.exists(path):
                continue
            sl = prs.slides.add_slide(blank)
            with Image.open(path) as im:
                iw, ih = im.size
            img_aspect = iw / float(ih)
            if img_aspect > slide_aspect:
                w = SW; h = int(SW / img_aspect)
            else:
                h = SH; w = int(SH * img_aspect)
            sl.shapes.add_picture(path, int((SW - w) / 2), int((SH - h) / 2), width=w, height=h)
        prs.save(output_pptx)
        print(f"PPTX (image slides) saved: {output_pptx}")
        return output_pptx

    # ---- Otherwise build a native editable deck (title + metrics + tables) ----
    s = prs.slides.add_slide(blank)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(2.4), SW - Inches(1.4), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = context.get('title', 'Analytics Report')
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = context.get('subtitle', '')
    p2.font.size = Pt(20)
    p2.font.color.rgb = GREY
    p2.alignment = PP_ALIGN.CENTER

    # ---- Overview metrics slide ----
    metrics = context.get('metrics') or []
    if metrics:
        s = prs.slides.add_slide(blank)
        title_bar(s, 'Overview')
        per_row = 4
        gap = Inches(0.3)
        left0 = Inches(0.5)
        card_w = int((SW - Inches(1) - gap * (per_row - 1)) / per_row)
        card_h = Inches(1.9)
        GREENISH = RGBColor(0x86, 0xEF, 0xAC)
        REDDISH = RGBColor(0xFE, 0xCA, 0xCA)
        for i, m in enumerate(metrics):
            label, value = m[0], m[1]
            change = m[2] if len(m) > 2 else None
            row = i // per_row
            col = i % per_row
            x = int(left0 + col * (card_w + gap))
            top = int(Inches(1.7) + row * (card_h + Inches(0.35)))
            shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, card_w, card_h)
            shp.fill.solid()
            shp.fill.fore_color.rgb = DARK
            shp.line.color.rgb = DARK
            tf = shp.text_frame
            tf.word_wrap = True
            pv = tf.paragraphs[0]
            pv.text = str(value)
            pv.font.size = Pt(28)
            pv.font.bold = True
            pv.font.color.rgb = LIME
            pv.alignment = PP_ALIGN.CENTER
            pl = tf.add_paragraph()
            pl.text = label
            pl.font.size = Pt(12)
            pl.font.color.rgb = WHITE
            pl.alignment = PP_ALIGN.CENTER
            if change is not None:
                pc = tf.add_paragraph()
                pc.text = ('▲ ' if change >= 0 else '▼ ') + str(abs(change)) + '% vs prev'
                pc.font.size = Pt(10)
                pc.font.bold = True
                pc.font.color.rgb = GREENISH if change >= 0 else REDDISH
                pc.alignment = PP_ALIGN.CENTER

    # ---- Native editable charts (infographics) ----
    charts = context.get('charts') or []
    if charts:
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        for ch in charts:
            cats = [str(c) for c in (ch.get('categories') or [])]
            vals = [float(v) for v in (ch.get('values') or [])]
            if not cats or not vals:
                continue
            s = prs.slides.add_slide(blank)
            title_bar(s, ch.get('title', ''))
            cd = CategoryChartData()
            cd.categories = cats
            cd.add_series('Value', vals)
            ctype = XL_CHART_TYPE.LINE if ch.get('kind') == 'line' else XL_CHART_TYPE.COLUMN_CLUSTERED
            gframe = s.shapes.add_chart(ctype, Inches(0.6), Inches(1.5),
                                        SW - Inches(1.2), SH - Inches(2.2), cd)
            try:
                gframe.chart.has_legend = False
            except Exception:
                pass

    # ---- One slide per table ----
    for t in context.get('tables', []):
        headers = t.get('headers') or []
        rows = (t.get('rows') or [])[:14]
        if not headers or not rows:
            continue
        s = prs.slides.add_slide(blank)
        title_bar(s, t.get('title', ''))
        nrows = len(rows) + 1
        ncols = len(headers)
        left = Inches(0.5)
        top = Inches(1.4)
        width = SW - Inches(1)
        height = min(SH - Inches(1.7), Inches(0.4) * nrows)
        table = s.shapes.add_table(nrows, ncols, left, top, width, int(height)).table
        for c, h in enumerate(headers):
            cell = table.cell(0, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            cp = cell.text_frame.paragraphs[0]
            cp.text = str(h)
            cp.font.size = Pt(11)
            cp.font.bold = True
            cp.font.color.rgb = WHITE
        for r, row in enumerate(rows, start=1):
            for c in range(ncols):
                cell = table.cell(r, c)
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
                cp = cell.text_frame.paragraphs[0]
                cp.text = str(row[c]) if c < len(row) else ''
                cp.font.size = Pt(10)
                cp.font.color.rgb = DARK

    prs.save(output_pptx)
    print(f"Editable PPTX saved: {output_pptx}")
    return output_pptx
